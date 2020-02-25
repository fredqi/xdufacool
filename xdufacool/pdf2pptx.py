#!/usr/bin/env python3
#
# Filename: pdf2pptx.py
# Author: Fred Qi
# Created: 2020-02-14 22:43:36(+0800)
#
# Last-Updated: 2020-02-25 10:27:48(+0800) [by Fred Qi]
#     Update #: 171
# 

# Commentary:
#
#
# 

# Change Log:
#
#
# 

import os
import argparse

from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


def pdf_slide_to_pptx(slide_file, output_folder, template=None, dpi=200):
    """
    To convert a pdf slide given by slide_file to a PPTX file.
    """

    assert os.path.exists(slide_file)
    _, pdf_name = os.path.split(slide_file)
    basename, ext = os.path.splitext(pdf_name)
    assert ext.lower() == '.pdf'
    
    images = convert_from_path(slide_file, dpi=dpi)
    
    image_folder = os.path.join(output_folder, '.' + basename)
    if not os.path.exists(image_folder):
        os.mkdir(image_folder)

    image_paths = []
    for idx, image in enumerate(images):
        # print(type(image), image.getbbox())
        image_path = os.path.join(image_folder, f'image{idx:02d}.jpg')
        image.save(image_path)
        image_paths.append(image_path)

    if template is None:
        prs = Presentation()
    else:
        prs = Presentation(template)
    blank_slide_layout = prs.slide_layouts[6]

    left = top = 0
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_path, left, top,
                                 height=prs.slide_height)

    pptx_file = os.path.join(output_folder, basename + '.pptx')
    prs.save(pptx_file)


def pdf2pptx():
    """Parsing command line arguments and make the pdf to pptx convertion."""
    
    parser = argparse.ArgumentParser(description='Convert a PDF slide to a PPTX file')
    parser.add_argument('--dpi', '-d', default=200, type=int,
                        help='The dpi of the converted image of a page in a PDF slide.')
    parser.add_argument('--template', '-t', default='', type=str,
                        help='The template PPTX file to create new presentations.')
    parser.add_argument('--output_folder', '-o', default='pptx', type=str,
                        help='The output folder of the PPTX file.')
    parser.add_argument('slide', type=str,
                        help='The filename of the PDF slide to be converted.')

    args = parser.parse_args()
    # print(args)
    options = {'output_folder': args.output_folder, 'dpi': args.dpi}
    if os.path.isfile(args.template):
        options['template'] = args.template
    pdf_slide_to_pptx(args.slide, **options)


if __name__ == '__main__':
    pdf2pptx()
    
# 
# pdf2pptx.py ends here
